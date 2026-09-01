# SPDX-License-Identifier: Apache-2.0
# Copyright 2026 Fabio Campolim
"""Export the course deck to PPTX: one full-bleed screenshot per slide (final
fragment state) with the lecturer notes as editable speaker notes.

Needs the optional tooling env (course/tools/requirements.txt: playwright +
python-pptx). Output defaults to course/build/<deck title>.pptx (gitignored).

Usage (from pythtb-skill/):
    python course/tools/build_pptx.py [--out FILE.pptx]
"""

import argparse
import os
import sys
import tempfile

__version__ = "1.4.1"

HERE = os.path.dirname(os.path.abspath(__file__))
COURSE = os.path.abspath(os.path.join(HERE, ".."))
INDEX = os.path.join(COURSE, "deck", "index.html")
OUT = os.path.join(COURSE, "build", "Tight-Binding Physics with PythTB.pptx")


def export(index=INDEX, out=OUT, quiet=False):
    from playwright.sync_api import sync_playwright   # optional dependencies
    from pptx import Presentation
    from pptx.util import Inches

    shots = tempfile.mkdtemp(prefix="pythtb-course-pptx-")
    slides = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        page.goto("file:///" + index.replace("\\", "/"))
        page.wait_for_function("() => window.Reveal && Reveal.isReady()")
        page.evaluate("Reveal.configure({controls:false, progress:false, transition:'none'})")
        page.evaluate("document.querySelectorAll('.deck-nav,.deck-progress,.deck-divider')"
                      ".forEach(e => e.style.display = 'none')")
        indices = page.evaluate(
            "Reveal.getSlides().map(s => { const i = Reveal.getIndices(s); return [i.h, i.v || 0]; })")
        for n, (h, v) in enumerate(indices):
            page.evaluate(f"Reveal.slide({h}, {v}, -1)")
            page.wait_for_timeout(200)
            while page.evaluate("Reveal.availableFragments().next === true"):
                page.evaluate("Reveal.nextFragment()")
                page.wait_for_timeout(80)
            page.wait_for_timeout(300)
            sid = page.evaluate("Reveal.getCurrentSlide().id") or f"slide-{n}"
            notes = page.evaluate("(id) => (window.DECK_CONTENT.slides[id] || {}).notes || ''", sid)
            png = os.path.join(shots, f"{n:03d}-{sid}.png")
            page.screenshot(path=png)
            slides.append((png, notes))
            if not quiet:
                print(f"  [{n + 1:2d}/{len(indices)}] {sid}")
        browser.close()

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    for png, notes in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
        slide.notes_slide.notes_text_frame.text = notes
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    if not quiet:
        print(f"ok: {out} ({len(slides)} slides)")
    return len(slides)


def build_parser():
    p = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    p.add_argument("--index", default=INDEX, help="deck to export")
    p.add_argument("--out", default=OUT, help="output .pptx")
    p.add_argument("-q", "--quiet", action="store_true")
    p.add_argument("--version", action="version", version=f"%(prog)s {__version__}")
    return p


def main(argv=None):
    a = build_parser().parse_args(argv)
    export(a.index, a.out, a.quiet)
    return 0


if __name__ == "__main__":
    sys.exit(main())
